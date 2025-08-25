import os
import uuid
import mimetypes
from pathlib import Path
from typing import List, Dict, Any, Optional
from fastapi import UploadFile, HTTPException, status
from docx import Document as DocxDocument
from pypdf import PdfReader
from openpyxl import load_workbook
from pptx import Presentation

# LlamaIndex imports - optional for now
try:
    from llama_index.core import SimpleDirectoryReader
    from llama_index.core.schema import Document as LlamaDocument
    from llama_index.core.node_parser import (
        SentenceSplitter,
        SemanticSplitterNodeParser,
        TokenTextSplitter,
        RecursiveCharacterTextSplitter
    )
    from llama_index.core.embeddings import BaseEmbedding
    from llama_index.core.text_splitter import SentenceSplitter as TextSentenceSplitter
    LLAMAINDEX_AVAILABLE = True
except ImportError:
    SimpleDirectoryReader = None
    LlamaDocument = None
    SentenceSplitter = None
    SemanticSplitterNodeParser = None
    TokenTextSplitter = None
    RecursiveCharacterTextSplitter = None
    BaseEmbedding = None
    TextSentenceSplitter = None
    LLAMAINDEX_AVAILABLE = False


class DocumentProcessor:
    ALLOWED_EXTENSIONS = {'.pdf', '.txt', '.md', '.docx', '.xlsx', '.xls', '.pptx'}
    MAX_FILE_SIZE = 30 * 1024 * 1024  # 30MB
    UPLOAD_DIR = "uploads"
    
    # LlamaIndex supported file types
    LLAMAINDEX_SUPPORTED = {'.pdf', '.txt', '.md', '.docx'}
    CUSTOM_PARSER_REQUIRED = {'.xlsx', '.xls', '.pptx'}
    
    # Chunking strategy configurations
    CHUNK_STRATEGIES = {
        'sentence': {
            'chunk_size': 512,
            'chunk_overlap': 50,
            'separator': None  # Uses intelligent sentence boundaries
        },
        'semantic': {
            'buffer_size': 1,
            'percentile_cutoff': 95
        },
        'token': {
            'chunk_size': 512,
            'chunk_overlap': 50
        },
        'recursive': {
            'chunk_size': 1000,
            'chunk_overlap': 100,
            'separators': ['\n\n', '\n', '。', '、', '.', ',', ' ', '']
        },
        'default': {
            'chunk_size': 1000,
            'overlap': 100
        }
    }

    @classmethod
    def validate_file(cls, file: UploadFile) -> None:
        if not file.filename:
            raise HTTPException(
                status_code=status.HTTP_400_BAD_REQUEST,
                detail="ファイル名が指定されていません"
            )

        # ファイル拡張子の確認
        file_extension = Path(file.filename).suffix.lower()
        if file_extension not in cls.ALLOWED_EXTENSIONS:
            raise HTTPException(
                status_code=status.HTTP_400_BAD_REQUEST,
                detail=f"サポートされていないファイル形式です。許可される拡張子: {', '.join(cls.ALLOWED_EXTENSIONS)}"
            )

        # ファイルサイズの確認 (content-lengthが利用可能な場合)
        if hasattr(file, 'size') and file.size and file.size > cls.MAX_FILE_SIZE:
            raise HTTPException(
                status_code=status.HTTP_413_REQUEST_ENTITY_TOO_LARGE,
                detail=f"ファイルサイズが上限（{cls.MAX_FILE_SIZE // (1024*1024)}MB）を超えています"
            )

    @classmethod
    async def save_file(cls, file: UploadFile) -> tuple[str, str, int]:
        # アップロードディレクトリの確保
        upload_path = Path(cls.UPLOAD_DIR)
        upload_path.mkdir(parents=True, exist_ok=True)

        # 一意なファイル名の生成
        file_extension = Path(file.filename).suffix.lower()
        unique_filename = f"{uuid.uuid4()}{file_extension}"
        file_path = upload_path / unique_filename

        # ファイルの保存とサイズ測定
        file_size = 0
        try:
            with open(file_path, "wb") as buffer:
                content = await file.read()
                file_size = len(content)
                
                # サイズ再確認
                if file_size > cls.MAX_FILE_SIZE:
                    raise HTTPException(
                        status_code=status.HTTP_413_REQUEST_ENTITY_TOO_LARGE,
                        detail=f"ファイルサイズが上限（{cls.MAX_FILE_SIZE // (1024*1024)}MB）を超えています"
                    )
                
                buffer.write(content)
        except Exception as e:
            # エラー時はファイルを削除
            if file_path.exists():
                file_path.unlink()
            raise HTTPException(
                status_code=status.HTTP_500_INTERNAL_SERVER_ERROR,
                detail=f"ファイルの保存中にエラーが発生しました: {str(e)}"
            )

        return str(file_path), unique_filename, file_size

    @classmethod
    def get_mime_type(cls, filename: str) -> str:
        mime_type, _ = mimetypes.guess_type(filename)
        return mime_type or 'application/octet-stream'

    @classmethod
    def delete_file(cls, file_path: str) -> None:
        try:
            if os.path.exists(file_path):
                os.remove(file_path)
        except Exception:
            pass  # ファイル削除エラーは無視

    @classmethod
    def read_text_content(cls, file_path: str, mime_type: str) -> str:
        try:
            file_ext = Path(file_path).suffix.lower()
            
            if mime_type in ['text/plain', 'text/markdown']:
                with open(file_path, 'r', encoding='utf-8') as f:
                    return f.read()
            elif mime_type == 'application/pdf' or file_ext == '.pdf':
                return cls._extract_pdf_text(file_path)
            elif (
                mime_type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
                or file_ext == ".docx"
            ):
                return cls._extract_docx_text(file_path)
            elif file_ext in ['.xlsx', '.xls']:
                return cls._extract_excel_text(file_path)
            elif file_ext == '.pptx':
                return cls._extract_pptx_text(file_path)
            else:
                raise ValueError(f"サポートされていないMIMEタイプ: {mime_type}")
        except Exception as e:
            raise HTTPException(
                status_code=status.HTTP_500_INTERNAL_SERVER_ERROR,
                detail=f"ファイルの内容読み取りエラー: {str(e)}"
            )

    @classmethod
    def _extract_pdf_text(cls, file_path: str) -> str:
        # PDF処理は今回の実装では簡略化
        # 実際の実装では pypdf や pdfplumber などのライブラリを使用
        try:
            reader = PdfReader(file_path)
            text_chunks = []
            for page in reader.pages:
                try:
                    text_chunks.append(page.extract_text() or "")
                except Exception:
                    # Ignore individual page errors to ensure best-effort extraction
                    continue
            return "\n".join(text_chunks)
        except Exception as e:
            raise RuntimeError(f"PDFテキスト抽出エラー: {str(e)}") from e

    @classmethod
    def _extract_docx_text(cls, file_path: str) -> str:
        """
        Extract text from a DOCX file by concatenating all paragraph texts.
        """
        try:
            doc = DocxDocument(file_path)
            paragraphs = [p.text for p in doc.paragraphs if p.text]
            return "\n".join(paragraphs)
        except Exception as e:
            raise RuntimeError(f"DOCXテキスト抽出エラー: {str(e)}") from e

    @classmethod
    def _extract_excel_text(cls, file_path: str) -> str:
        """
        Extract text from Excel files (.xlsx, .xls)
        """
        try:
            workbook = load_workbook(filename=file_path, data_only=True)
            text_parts = []
            
            for sheet_name in workbook.sheetnames:
                worksheet = workbook[sheet_name]
                sheet_text = [f"[Sheet: {sheet_name}]"]
                
                for row in worksheet.iter_rows(values_only=True):
                    row_text = [str(cell) if cell is not None else "" for cell in row]
                    if any(cell.strip() for cell in row_text):
                        sheet_text.append("\t".join(row_text))
                
                if len(sheet_text) > 1:  # Has content beyond the sheet name
                    text_parts.extend(sheet_text)
                    text_parts.append("")  # Add separator between sheets
            
            return "\n".join(text_parts)
        except Exception as e:
            raise RuntimeError(f"Excelテキスト抽出エラー: {str(e)}") from e
    
    @classmethod
    def _extract_pptx_text(cls, file_path: str) -> str:
        """
        Extract text from PowerPoint files (.pptx)
        """
        try:
            presentation = Presentation(file_path)
            text_parts = []
            
            for slide_num, slide in enumerate(presentation.slides, 1):
                slide_text = [f"[Slide {slide_num}]"]
                
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text.append(shape.text.strip())
                
                if len(slide_text) > 1:  # Has content beyond slide number
                    text_parts.extend(slide_text)
                    text_parts.append("")  # Add separator between slides
            
            return "\n".join(text_parts)
        except Exception as e:
            raise RuntimeError(f"PowerPointテキスト抽出エラー: {str(e)}") from e
    
    @classmethod
    def create_llamaindex_documents(
        cls, file_paths: List[str], metadata: Optional[Dict[str, Any]] = None
    ) -> List[Any]:
        """
        Create LlamaIndex documents from file paths using SimpleDirectoryReader
        and custom parsers for unsupported formats.
        Returns regular dicts if LlamaIndex is not available.
        """
        documents = []
        
        if not LLAMAINDEX_AVAILABLE:
            # Fallback to basic document processing
            for file_path in file_paths:
                try:
                    mime_type = cls.get_mime_type(file_path)
                    content = cls.read_text_content(file_path, mime_type)
                    doc_metadata = {"file_path": file_path, "file_name": Path(file_path).name}
                    if metadata:
                        doc_metadata.update(metadata)
                    documents.append({"text": content, "metadata": doc_metadata})
                except Exception as e:
                    print(f"Document processing failed for {file_path}: {str(e)}")
            return documents
        llamaindex_files = []
        custom_parser_files = []
        
        # Separate files by processing method
        for file_path in file_paths:
            file_ext = Path(file_path).suffix.lower()
            if file_ext in cls.LLAMAINDEX_SUPPORTED:
                llamaindex_files.append(file_path)
            elif file_ext in cls.CUSTOM_PARSER_REQUIRED:
                custom_parser_files.append(file_path)
        
        # Process files supported by LlamaIndex
        if llamaindex_files:
            try:
                reader = SimpleDirectoryReader(input_files=llamaindex_files)
                documents.extend(reader.load_data())
            except Exception as e:
                print(f"LlamaIndex processing error: {str(e)}")
                # Fallback to custom processing
                for file_path in llamaindex_files:
                    try:
                        mime_type = cls.get_mime_type(file_path)
                        content = cls.read_text_content(file_path, mime_type)
                        doc_metadata = {"file_path": file_path, "file_name": Path(file_path).name}
                        if metadata:
                            doc_metadata.update(metadata)
                        if LLAMAINDEX_AVAILABLE:
                            documents.append(LlamaDocument(text=content, metadata=doc_metadata))
                        else:
                            documents.append({"text": content, "metadata": doc_metadata})
                    except Exception as fallback_error:
                        print(f"Fallback processing failed for {file_path}: {fallback_error}")
        
        # Process files requiring custom parsers
        for file_path in custom_parser_files:
            try:
                mime_type = cls.get_mime_type(file_path)
                content = cls.read_text_content(file_path, mime_type)
                doc_metadata = {"file_path": file_path, "file_name": Path(file_path).name}
                if metadata:
                    doc_metadata.update(metadata)
                if LLAMAINDEX_AVAILABLE:
                    documents.append(LlamaDocument(text=content, metadata=doc_metadata))
                else:
                    documents.append({"text": content, "metadata": doc_metadata})
            except Exception as e:
                print(f"Custom parser failed for {file_path}: {str(e)}")
        
        return documents
    
    @classmethod
    def chunk_text(
        cls, 
        content: str, 
        chunk_size: int = 1000, 
        overlap: int = 100,
        strategy: str = 'default',
        project_type: str = None,
        embedding_model = None,
        **kwargs  # Additional project-specific configuration
    ) -> List[str]:
        """Intelligent chunking with multiple strategies"""
        if not content.strip():
            return []
            
        # Use intelligent chunking if LlamaIndex is available
        if LLAMAINDEX_AVAILABLE and strategy != 'default':
            return cls._intelligent_chunk_text(content, strategy, project_type, embedding_model)
        
        # Fallback to original chunking method
        chunks = []
        start = 0
        content_length = len(content)

        while start < content_length:
            end = min(start + chunk_size, content_length)
            
            # 文の区切りを探す
            if end < content_length:
                # 改行、句読点で区切る
                for delimiter in ['\n\n', '\n', '。', '.']:
                    last_delimiter = content.rfind(delimiter, start, end)
                    if last_delimiter > start:
                        end = last_delimiter + len(delimiter)
                        break

            chunk = content[start:end].strip()
            if chunk:
                chunks.append(chunk)

            start = max(start + 1, end - overlap)

        return chunks
    
    @classmethod
    def _intelligent_chunk_text(
        cls,
        content: str,
        strategy: str,
        project_type: str = None,
        embedding_model = None
    ) -> List[str]:
        """Advanced chunking using LlamaIndex node parsers"""
        if not LLAMAINDEX_AVAILABLE:
            return cls.chunk_text(content)  # Fallback
            
        try:
            # Create a LlamaDocument for processing
            document = LlamaDocument(text=content)
            
            # Select and configure node parser based on strategy
            node_parser = cls._get_node_parser(strategy, project_type, embedding_model)
            
            # Parse document into nodes
            nodes = node_parser.get_nodes_from_documents([document])
            
            # Extract text content from nodes
            chunks = [node.text for node in nodes if node.text.strip()]
            
            return chunks
            
        except Exception as e:
            print(f"Intelligent chunking failed: {str(e)}. Falling back to default.")
            return cls.chunk_text(content)
    
    @classmethod
    def _get_node_parser(cls, strategy: str, project_type: str = None, embedding_model = None):
        """Get appropriate node parser based on strategy and project type"""
        config = cls.CHUNK_STRATEGIES.get(strategy, cls.CHUNK_STRATEGIES['default'])
        
        if strategy == 'sentence':
            return SentenceSplitter(
                chunk_size=config['chunk_size'],
                chunk_overlap=config['chunk_overlap']
            )
        elif strategy == 'semantic' and embedding_model is not None:
            return SemanticSplitterNodeParser(
                buffer_size=config['buffer_size'],
                percentile_cutoff=config['percentile_cutoff'],
                embed_model=embedding_model
            )
        elif strategy == 'token':
            return TokenTextSplitter(
                chunk_size=config['chunk_size'],
                chunk_overlap=config['chunk_overlap']
            )
        elif strategy == 'recursive':
            return RecursiveCharacterTextSplitter(
                chunk_size=config['chunk_size'],
                chunk_overlap=config['chunk_overlap'],
                separators=config['separators']
            )
        else:
            # Default to sentence splitter for unknown strategies
            return SentenceSplitter(
                chunk_size=512,
                chunk_overlap=50
            )
    
    @classmethod
    def get_optimal_chunking_strategy(cls, project_type: str = None, file_types: List[str] = None) -> str:
        """Determine optimal chunking strategy based on project and file types"""
        # Project-specific optimization
        if project_type:
            project_type = project_type.lower()
            if 'research' in project_type or 'academic' in project_type:
                return 'semantic'  # Better for academic papers
            elif 'legal' in project_type or 'contract' in project_type:
                return 'sentence'  # Preserve legal structure
            elif 'technical' in project_type or 'manual' in project_type:
                return 'recursive'  # Better for structured documents
            elif 'code' in project_type or 'development' in project_type:
                return 'token'  # Better for code documentation
        
        # File type-based optimization
        if file_types:
            pdf_count = sum(1 for ft in file_types if ft.lower() == '.pdf')
            md_count = sum(1 for ft in file_types if ft.lower() == '.md')
            docx_count = sum(1 for ft in file_types if ft.lower() == '.docx')
            
            if pdf_count > len(file_types) * 0.5:  # Majority PDFs
                return 'sentence'  # PDFs often have complex formatting
            elif md_count > len(file_types) * 0.3:  # Significant markdown
                return 'recursive'  # Markdown has clear hierarchical structure
            elif docx_count > len(file_types) * 0.4:  # Majority Word docs
                return 'sentence'  # Preserve document structure
        
        return 'sentence'  # Safe default for most use cases